#!/usr/bin/env python3
"""
Build steering-committee-deck.pptx — tuned for Google Slides import.

Typography: **Montserrat** (slide titles / kickers) + **Roboto** (body, tables, bullets).
Both are on the Google Slides font list and read closer to modern web decks than Arial.

Layout: explicit table column widths (~58% label / ~42% values) and row heights so
12–13pt text fits without overlapping the footer. Principles: one slide, 4+3 compact cards.

Cover slide: light blue background + navy type (aligned with HTML); decorative ovals on the right.

Requires: python-pptx (see requirements-steering-deck.txt).
"""
from __future__ import annotations

from datetime import date
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

ROOT = Path(__file__).resolve().parent
OUT = ROOT / "steering-committee-deck.pptx"

# Google Slides–friendly stack (HTML deck uses Plus Jakarta Sans; closest native pair)
FONT_TITLE = "Montserrat"
FONT_BODY = "Roboto"

BG = RGBColor(0, 11, 25)
WHITE = RGBColor(232, 236, 240)
COVER_BG = RGBColor(232, 242, 252)  # #e8f2fc — light cover (matches HTML)
COVER_NAVY = RGBColor(15, 31, 61)  # primary text on light cover
COVER_MUTED = RGBColor(71, 85, 105)  # slate-600
COVER_ACCENT = RGBColor(30, 64, 175)  # deep blue accent
MUTED = RGBColor(156, 163, 175)
ACCENT = RGBColor(61, 127, 224)
BR = RGBColor(34, 197, 94)
BR_BRIGHT = RGBColor(74, 222, 128)  # slightly brighter for table headers on import
AR = RGBColor(59, 130, 246)
AR_BRIGHT = RGBColor(96, 165, 250)
MX = RGBColor(239, 68, 68)
STONE = RGBColor(245, 158, 11)
DEV = RGBColor(167, 139, 250)
LAUNCH_BR = RGBColor(34, 197, 94)
LAUNCH_AR = RGBColor(59, 130, 246)

PANEL_BG = RGBColor(28, 32, 48)
PANEL_LINE = RGBColor(90, 100, 128)

TBL_HDR = RGBColor(0, 55, 110)
TBL_ROW_A = RGBColor(12, 28, 48)
TBL_ROW_B = RGBColor(18, 34, 54)

# Sizes tuned for Slides (avoid <10pt; tables 12–13pt; cover/principles ≥11pt for import)
PT_KICKER = 12
PT_TITLE = 30
PT_SUBTITLE = 12
PT_PRINCIPLES_SUB = 14  # principles slide sub-header (dense paragraph)
PT_CARD_NAME = 20
PT_CARD_LEDE = 13
PT_CARD_BULLET = 14
PT_INSIGHTS = 11
PT_TABLE = 13
PT_FOOT = 10
PT_GANTT = 11
PT_PRINCIPLE = 14
PT_BOX_HEAD = 11
PT_BOX_BODY = 12
# Principle cards — keep ≥11pt so Google Slides does not render microscopic text
PT_PRIN_COMPACT_HEAD = 12
PT_PRIN_COMPACT_BODY = 11


def _font_title(paragraph) -> None:
    paragraph.font.name = FONT_TITLE


def _font_body(paragraph) -> None:
    paragraph.font.name = FONT_BODY


def _set_slide_bg(slide, color: RGBColor) -> None:
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = color


def _add_title(slide, text: str, top: float, height: float, size: int = PT_TITLE) -> None:
    box = slide.shapes.add_textbox(Inches(0.55), Inches(top), Inches(12.25), Inches(height))
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.bold = True
    p.font.color.rgb = WHITE
    _font_title(p)


def _add_kicker(slide, text: str, top: float) -> None:
    box = slide.shapes.add_textbox(Inches(0.55), Inches(top), Inches(12.0), Inches(0.4))
    p = box.text_frame.paragraphs[0]
    p.text = text
    p.font.size = Pt(PT_KICKER)
    p.font.bold = True
    p.font.color.rgb = ACCENT
    _font_title(p)


def _body_box(
    slide,
    left: float,
    top: float,
    width: float,
    height: float,
    text: str,
    size: int,
    color: RGBColor = WHITE,
    font: str = "body",
) -> None:
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.TOP
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.color.rgb = color
    if font == "title":
        _font_title(p)
    else:
        _font_body(p)
    p.line_spacing = 1.2


def _cell_fill(cell, rgb: RGBColor) -> None:
    try:
        cell.fill.solid()
        cell.fill.fore_color.rgb = rgb
    except (AttributeError, TypeError):
        pass


def _table_apply_column_widths(table, total_width_inches: float) -> None:
    """~58% metrics / labels, ~42% values — balances wide number columns (Slides)."""
    w0 = int(Inches(total_width_inches * 0.58))
    w1 = int(Inches(total_width_inches * 0.42))
    table.columns[0].width = w0
    table.columns[1].width = w1


def _table_apply_row_heights(table, height_pt: int) -> None:
    for row in table.rows:
        row.height = Pt(height_pt)


def _table_from_rows(
    slide,
    left: float,
    top: float,
    width_inches: float,
    height_inches: float,
    rows: list[tuple[str, str]],
    accent_rgb: RGBColor,
    font_pt: int = PT_TABLE,
    row_height_pt: int = 21,
) -> None:
    if not rows:
        return
    rcount = len(rows)
    graphic_frame = slide.shapes.add_table(
        rcount, 2, Inches(left), Inches(top), Inches(width_inches), Inches(height_inches)
    )
    table = graphic_frame.table
    _table_apply_column_widths(table, width_inches)

    for i, (a, b) in enumerate(rows):
        c0 = table.cell(i, 0)
        c1 = table.cell(i, 1)
        c0.text = a
        c1.text = b
        is_sep = a.strip() == "—" and b.strip() == "—"
        is_section_title = b == "" and a and not is_sep
        row_fill = TBL_HDR if (i == 0 or is_section_title or is_sep) else (TBL_ROW_A if i % 2 == 0 else TBL_ROW_B)
        if is_sep:
            row_fill = RGBColor(25, 40, 62)

        for ci, cell in enumerate((c0, c1)):
            _cell_fill(cell, row_fill)
            for para in cell.text_frame.paragraphs:
                para.font.size = Pt(font_pt)
                _font_body(para)
                para.font.color.rgb = WHITE
            cell.text_frame.paragraphs[0].alignment = PP_ALIGN.LEFT if ci == 0 else PP_ALIGN.RIGHT
            if is_section_title or i == 0:
                cell.text_frame.paragraphs[0].font.bold = True
                if ci == 0 or (b.strip() and ci == 1):
                    cell.text_frame.paragraphs[0].font.color.rgb = accent_rgb
                else:
                    cell.text_frame.paragraphs[0].font.color.rgb = MUTED
            elif is_sep:
                cell.text_frame.paragraphs[0].font.color.rgb = MUTED

    _table_apply_row_heights(table, row_height_pt)


def _add_cover_slide(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide, COVER_BG)

    # Decorative blobs (right side) — flat organic shapes
    blob_specs: list[tuple[float, float, float, float, RGBColor]] = [
        (7.1, 0.55, 1.55, 1.55, RGBColor(180, 210, 245)),
        (9.5, 1.15, 1.2, 1.35, RGBColor(200, 225, 250)),
        (6.8, 4.2, 1.85, 1.55, RGBColor(165, 200, 235)),
        (10.2, 4.9, 1.2, 1.1, RGBColor(210, 230, 252)),
    ]
    for lx, ty, w, h, rgb in blob_specs:
        o = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(lx), Inches(ty), Inches(w), Inches(h))
        o.fill.solid()
        o.fill.fore_color.rgb = rgb
        o.line.fill.background()

    box = slide.shapes.add_textbox(Inches(0.72), Inches(1.05), Inches(6.85), Inches(5.55))
    tf = box.text_frame
    tf.word_wrap = True
    p0 = tf.paragraphs[0]
    p0.text = "POS + PAYMENTS · PRICING STRATEGY"
    p0.font.size = Pt(13)
    p0.font.bold = True
    p0.font.color.rgb = COVER_ACCENT
    p0.alignment = PP_ALIGN.LEFT
    _font_title(p0)

    p2 = tf.add_paragraph()
    p2.text = "Steering committee deck"
    p2.font.size = Pt(38)
    p2.font.bold = True
    p2.font.color.rgb = COVER_NAVY
    p2.alignment = PP_ALIGN.LEFT
    _font_title(p2)
    p2.space_before = Pt(10)

    p3 = tf.add_paragraph()
    p3.text = "First readout — offline payments, PDV, and commercial roadmap"
    p3.font.size = Pt(16)
    p3.font.color.rgb = COVER_NAVY
    p3.alignment = PP_ALIGN.LEFT
    _font_body(p3)
    p3.space_before = Pt(8)

    pobj = tf.add_paragraph()
    pobj.text = (
        "Objectives: align expectations on next steps and long-term view; discuss pricing principles"
    )
    pobj.font.size = Pt(15)
    pobj.font.bold = True
    pobj.font.color.rgb = COVER_NAVY
    pobj.alignment = PP_ALIGN.LEFT
    _font_body(pobj)
    pobj.space_before = Pt(14)
    pobj.line_spacing = 1.25

    pdt = tf.add_paragraph()
    pdt.text = "March 2026"
    pdt.font.size = Pt(13)
    pdt.font.color.rgb = COVER_MUTED
    pdt.alignment = PP_ALIGN.LEFT
    _font_body(pdt)
    pdt.space_before = Pt(18)

    p5 = tf.add_paragraph()
    p5.text = "Designed HTML deck (GitHub Pages)"
    p5.font.size = Pt(12)
    p5.font.bold = True
    p5.font.color.rgb = COVER_ACCENT
    p5.alignment = PP_ALIGN.LEFT
    _font_body(p5)
    p5.space_before = Pt(22)

    p6 = tf.add_paragraph()
    p6.text = "https://henryjbpetry.github.io/pos_payments/pos-pricing-model/steering-committee-deck.html"
    p6.font.size = Pt(11)
    p6.font.color.rgb = COVER_NAVY
    p6.alignment = PP_ALIGN.LEFT
    _font_body(p6)

    p8 = tf.add_paragraph()
    p8.text = "Source file"
    p8.font.size = Pt(12)
    p8.font.bold = True
    p8.font.color.rgb = COVER_MUTED
    p8.alignment = PP_ALIGN.LEFT
    _font_body(p8)
    p8.space_before = Pt(10)

    p9 = tf.add_paragraph()
    p9.text = "https://github.com/henryjbpetry/pos_payments/blob/main/pos-pricing-model/steering-committee-deck.html"
    p9.font.size = Pt(11)
    p9.font.color.rgb = COVER_ACCENT
    p9.alignment = PP_ALIGN.LEFT
    _font_body(p9)


def _principle_card(
    slide,
    left: float,
    top: float,
    w: float,
    h: float,
    num: str,
    icon: str,
    title: str,
    why: str,
    *,
    title_pt: int = PT_BOX_HEAD + 1,
    body_pt: int = PT_BOX_BODY,
    line_before_body: int = 8,
) -> None:
    shp = slide.shapes.add_shape(1, Inches(left), Inches(top), Inches(w), Inches(h))
    shp.fill.solid()
    shp.fill.fore_color.rgb = RGBColor(22, 30, 52)
    shp.line.color.rgb = PANEL_LINE
    shp.line.width = Pt(1)
    tb = slide.shapes.add_textbox(Inches(left + 0.1), Inches(top + 0.07), Inches(w - 0.2), Inches(h - 0.12))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.TOP
    tf.margin_bottom = Inches(0)
    tf.margin_top = Inches(0)
    p0 = tf.paragraphs[0]
    p0.text = f"{num}  {icon}  {title}"
    p0.font.size = Pt(title_pt)
    p0.font.bold = True
    p0.font.color.rgb = WHITE
    p0.line_spacing = 1.12
    _font_title(p0)
    p1 = tf.add_paragraph()
    p1.text = why
    p1.font.size = Pt(body_pt)
    p1.font.color.rgb = MUTED
    _font_body(p1)
    p1.space_before = Pt(line_before_body)
    p1.line_spacing = 1.15


def build() -> None:
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    _add_cover_slide(prs)

    # --- Slide 0: Where we are ---
    slide0 = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide0, BG)
    _add_kicker(slide0, "POS · PRICING STRATEGY · STEERING COMMITTEE", 0.35)
    _add_title(slide0, "Where we are — by country", 0.58, 0.55, PT_TITLE)

    cards = [
        (
            "Brazil",
            "🇧🇷",
            "~99% toward Stone — pricing is the next big deliverable.",
            BR,
            (
                "Development: prioritizing work that overlaps multi-acquirer.",
                "Partner: ~99% Stone — tech simplicity, logistics.",
                "Next: pricing as the major milestone.",
                "Dependencies: Stone negotiations, joint tech roadmap.",
            ),
        ),
        (
            "Argentina",
            "🇦🇷",
            "PSP-first — QR interoperable only if PDV can build it.",
            AR,
            (
                "Pagonube team fully on PSP development.",
                "QR interoperable path under analysis.",
                "After PSP nears launch: revisit acquirer terminals.",
            ),
        ),
        (
            "Mexico",
            "🇲🇽",
            "Capacity-gated — two devs to Pagonube; Stripe engaged.",
            MX,
            (
                "Blocked on two new developers ceded to Pagonube.",
                "Stripe stays engaged while decision is open.",
            ),
        ),
    ]
    cw = 3.85
    gap = 0.35
    left0 = 0.55
    for i, (name, flag, sub, col, bullets) in enumerate(cards):
        x = left0 + i * (cw + gap)
        shp = slide0.shapes.add_shape(1, Inches(x), Inches(1.32), Inches(cw), Inches(5.35))
        shp.fill.solid()
        shp.fill.fore_color.rgb = PANEL_BG
        shp.line.color.rgb = PANEL_LINE
        shp.line.width = Pt(1)

        tb = slide0.shapes.add_textbox(Inches(x + 0.18), Inches(1.42), Inches(cw - 0.36), Inches(0.55))
        p = tb.text_frame.paragraphs[0]
        p.text = f"{flag}  {name}"
        p.font.size = Pt(PT_CARD_NAME)
        p.font.bold = True
        p.font.color.rgb = col
        _font_title(p)

        tb2 = slide0.shapes.add_textbox(Inches(x + 0.18), Inches(2.05), Inches(cw - 0.36), Inches(0.95))
        p2 = tb2.text_frame.paragraphs[0]
        p2.text = sub
        p2.font.size = Pt(PT_CARD_LEDE)
        p2.font.bold = True
        p2.font.color.rgb = WHITE
        _font_body(p2)

        tb3 = slide0.shapes.add_textbox(Inches(x + 0.18), Inches(3.05), Inches(cw - 0.36), Inches(3.45))
        tf3 = tb3.text_frame
        tf3.word_wrap = True
        for j, b in enumerate(bullets):
            pr = tf3.paragraphs[0] if j == 0 else tf3.add_paragraph()
            pr.text = f"• {b}"
            pr.font.size = Pt(PT_CARD_BULLET)
            pr.font.color.rgb = WHITE
            _font_body(pr)
            pr.line_spacing = 1.35
            pr.space_after = Pt(6)

    foot = slide0.shapes.add_textbox(Inches(0.55), Inches(6.92), Inches(12.0), Inches(0.42))
    fp = foot.text_frame.paragraphs[0]
    fp.text = "Internal steering snapshot — Mar 2026."
    fp.font.size = Pt(PT_FOOT)
    fp.font.color.rgb = MUTED
    _font_body(fp)

    # --- Economics 1/2: Brazil (compact header + table with safe footer gap) ---
    slide_e1 = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide_e1, BG)
    _add_kicker(slide_e1, "WHERE WE'RE GOING · UPDATED BUSINESS CASE (1/2)", 0.28)
    _add_title(
        slide_e1,
        "The economics justify a sales-led bet — if we sell the stack, not PoS alone",
        0.5,
        0.68,
        26,
    )
    sub = slide_e1.shapes.add_textbox(Inches(0.55), Inches(1.08), Inches(12.15), Inches(0.5))
    sp = sub.text_frame.paragraphs[0]
    sp.text = "Modeled unit economics — monthly USD per AE · This slide: Brazil (Stone). Next: Argentina (dLocal)."
    sp.font.size = Pt(PT_SUBTITLE)
    sp.font.color.rgb = MUTED
    _font_body(sp)

    insights = (
        "Key insights · PoS + offline pay → ~2–3 mo payback vs ~8 mo PoS-only · "
        "UpMarket = fewer deals, higher $/mo · MX out of scope until FX/local economics lock."
    )
    _body_box(slide_e1, 0.55, 1.55, 12.2, 0.42, insights, PT_INSIGHTS, MUTED)

    br_rows: list[tuple[str, str]] = [
        ("Brazil — Stone", ""),
        ("PoS — SaaS assumptions", ""),
        ("PoS take-rate (% of GMV)", "0.25%"),
        ("PoS subscription (USD/mo)", "13.1"),
        ("Gross margin PoS", "75.0%"),
        ("Upper SMB (US$5k–20k/mo, avg GMV US$9,963)", ""),
        ("Payback total / PoS-only", "2.9 mo / 7.6 mo"),
        ("Deals / month", "15"),
        ("Revenue total", "US$4,964"),
        ("Gross margin total", "US$1,139"),
        ("Commercial cost", "US$3,248"),
        ("UpMarket (US$20k–200k/mo, avg GMV US$36,917)", ""),
        ("Payback total / PoS-only", "2.8 mo / 8.8 mo"),
        ("Deals / month", "7"),
        ("Revenue total", "US$8,336"),
        ("Gross margin total", "US$1,720"),
        ("Commercial cost", "US$4,872"),
    ]
    # Table: start below insights; height sized for 19 rows × ~21pt + footer clearance
    _table_from_rows(
        slide_e1,
        0.55,
        2.05,
        12.2,
        4.72,
        br_rows,
        BR_BRIGHT,
        font_pt=PT_TABLE,
        row_height_pt=20,
    )

    foot1 = slide_e1.shapes.add_textbox(Inches(0.55), Inches(6.88), Inches(12.2), Inches(0.48))
    f1p = foot1.text_frame.paragraphs[0]
    f1p.text = (
        "Source: PoS GTM Sales Investment · Commercial cost 25% of monthly revenue total · Offline pay adoption 50%."
    )
    f1p.font.size = Pt(PT_FOOT)
    f1p.font.color.rgb = MUTED
    _font_body(f1p)

    # --- Economics 2/2: Argentina ---
    slide_e2 = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide_e2, BG)
    _add_kicker(slide_e2, "WHERE WE'RE GOING · UPDATED BUSINESS CASE (2/2)", 0.28)
    _add_title(slide_e2, "Argentina — same model, dLocal economics", 0.5, 0.52, 28)
    sub2 = slide_e2.shapes.add_textbox(Inches(0.55), Inches(1.05), Inches(12.15), Inches(0.48))
    s2 = sub2.text_frame.paragraphs[0]
    s2.text = "Monthly USD per AE · Upper SMB vs UpMarket (model export)."
    s2.font.size = Pt(PT_SUBTITLE)
    s2.font.color.rgb = MUTED
    _font_body(s2)

    ar_rows: list[tuple[str, str]] = [
        ("Argentina — dLocal", ""),
        ("PoS — SaaS assumptions", ""),
        ("PoS take-rate (% of GMV)", "0.25%"),
        ("PoS subscription (USD/mo)", "22.0"),
        ("Gross margin PoS", "75.0%"),
        ("Upper SMB (US$5k–20k/mo, avg GMV US$7,807)", ""),
        ("Payback total / PoS-only", "2.7 mo / 8.2 mo"),
        ("Deals / month", "15"),
        ("Revenue total", "US$6,386"),
        ("Gross margin total", "US$1,399"),
        ("Commercial cost", "US$3,816"),
        ("UpMarket (US$20k–200k/mo, avg GMV US$41,282)", ""),
        ("Payback total / PoS-only", "2.1 mo / 8.7 mo"),
        ("Deals / month", "7"),
        ("Revenue total", "US$15,099"),
        ("Gross margin total", "US$2,792"),
        ("Commercial cost", "US$5,724"),
    ]
    _table_from_rows(
        slide_e2,
        0.55,
        1.58,
        12.2,
        5.15,
        ar_rows,
        AR_BRIGHT,
        font_pt=PT_TABLE,
        row_height_pt=20,
    )

    foot2 = slide_e2.shapes.add_textbox(Inches(0.55), Inches(6.88), Inches(12.2), Inches(0.48))
    f2p = foot2.text_frame.paragraphs[0]
    f2p.text = "Reference: gopedroso.github.io/PoS/GTM-Sales-investment/ · Slide 3 · Sales-Led economics"
    f2p.font.size = Pt(PT_FOOT)
    f2p.font.color.rgb = ACCENT
    _font_body(f2p)

    # --- Gantt ---
    slide_g = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide_g, BG)
    _add_kicker(slide_g, "HOW WE'RE DOING IT", 0.26)
    _add_title(slide_g, "2026 roadmap — GTM & launch (Brazil · Argentina)", 0.48, 0.46, 26)

    d0 = date(2026, 3, 1)
    d1 = date(2026, 12, 31)
    total_days = (d1 - d0).days
    label_left, label_w = 0.5, 3.75
    gantt_left = label_left + label_w + 0.1
    gantt_w = 13.333 - gantt_left - 0.52
    chart_top = 1.05
    chart_h = 5.65
    axis_y = chart_top + chart_h + 0.06

    months = [
        (3, "Mar"),
        (4, "Apr"),
        (5, "May"),
        (6, "Jun"),
        (7, "Jul"),
        (8, "Aug"),
        (9, "Sep"),
        (10, "Oct"),
        (11, "Nov"),
        (12, "Dec"),
    ]
    for mnum, mlab in months:
        md = date(2026, mnum, 1)
        x = gantt_left + (md - d0).days / total_days * gantt_w
        t = slide_g.shapes.add_textbox(Inches(x), Inches(axis_y), Inches(0.56), Inches(0.26))
        tp = t.text_frame.paragraphs[0]
        tp.text = mlab
        tp.font.size = Pt(PT_GANTT)
        tp.font.color.rgb = MUTED
        tp.font.bold = True
        _font_body(tp)

    rows_data: list[tuple[str, date, date, RGBColor]] = [
        ("Credenciamento Stone", date(2026, 3, 28), date(2026, 5, 12), STONE),
        ("Development · BR *", date(2026, 4, 6), date(2026, 10, 15), DEV),
        ("Development · AR *", date(2026, 4, 6), date(2026, 8, 1), DEV),
        ("GTM · BR · Pricing strategy definition", date(2026, 3, 24), date(2026, 4, 13), BR),
        ("GTM · BR · New prices testing (no payments)", date(2026, 5, 1), date(2026, 6, 30), BR),
        ("GTM · BR · PM interviews & discovery", date(2026, 3, 24), date(2026, 5, 1), BR),
        ("GTM · BR · GTM launch planning", date(2026, 5, 20), date(2026, 7, 20), BR),
        ("GTM · BR · PIX campaign", date(2026, 6, 20), date(2026, 7, 30), BR),
        ("GTM · BR · Sales enablement (Payments + PdV)", date(2026, 7, 30), date(2026, 8, 30), BR),
        ("GTM · AR · Pricing strategy (QR)", date(2026, 4, 15), date(2026, 4, 20), AR),
        ("GTM · AR · PM & launch strategy (QR)", date(2026, 5, 20), date(2026, 6, 20), AR),
        ("Launch · BR · Pilot launch", date(2026, 9, 15), date(2026, 10, 15), LAUNCH_BR),
        ("Launch · BR · Post-launch review", date(2026, 10, 15), date(2026, 12, 11), LAUNCH_BR),
        ("Launch · AR · Pilot launch (QR)", date(2026, 7, 1), date(2026, 8, 1), LAUNCH_AR),
        ("Launch · AR · Post-launch review (QR)", date(2026, 8, 2), date(2026, 10, 15), LAUNCH_AR),
    ]

    row_h = chart_h / len(rows_data)
    for i, (label, sd, ed, col) in enumerate(rows_data):
        top_i = chart_top + i * row_h + 0.01
        left_pct = (sd - d0).days / total_days
        width_pct = max(0.004, (ed - sd).days / total_days)
        bar_left = gantt_left + left_pct * gantt_w
        bar_w = max(0.05, width_pct * gantt_w)

        lab = slide_g.shapes.add_textbox(Inches(label_left), Inches(top_i), Inches(label_w), Inches(max(0.24, row_h - 0.02)))
        lp = lab.text_frame.paragraphs[0]
        lp.text = label
        lp.font.size = Pt(PT_GANTT)
        lp.font.color.rgb = WHITE
        _font_body(lp)

        bar = slide_g.shapes.add_shape(
            1, Inches(bar_left), Inches(top_i + 0.02), Inches(bar_w), Inches(max(0.15, row_h - 0.08))
        )
        bar.fill.solid()
        bar.fill.fore_color.rgb = col
        bar.line.fill.background()

    leg = slide_g.shapes.add_textbox(Inches(0.5), Inches(6.86), Inches(12.35), Inches(0.55))
    p0 = leg.text_frame.paragraphs[0]
    p0.text = (
        "Legend: Stone · Development · GTM/Launch. * Development (6 Apr → each launch): macro bar — "
        "product breaks down internal milestones."
    )
    p0.font.size = Pt(PT_FOOT)
    p0.font.color.rgb = MUTED
    _font_body(p0)

    # --- Principles: single slide, 4 + 3 tight grid (matches HTML cards) ---
    principles: list[tuple[str, str, str, str]] = [
        ("1", "🏪", "Maximize offline channel results", "Measure in-store GMV & PDV-attributed revenue — not ecommerce."),
        ("2", "✨", "Radical simplicity for merchants", "Fewer variables; bill readable in one breath."),
        ("3", "🎯", "Align price communication to ICP reality", "Same list, different framing per archetype (PMM)."),
        ("4", "📊", "Ground list prices in non-promotional market bands", "Sustainable benchmarks; promos distort comparability."),
        ("5", "⏱️", "Respect settlement horizon sensitivity", "Prazo/anticipation ↔ effective MDR; surface in CRM."),
        ("6", "📈", "Monitor repricing and conversion jointly", "Price = product surface; watch funnel & tickets by segment."),
        ("7", "🔗", "Bundle economics with PoS", "Payments on physical GMV + PoS SaaS — not either line alone."),
    ]

    slide_p1 = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide_p1, BG)
    _add_kicker(slide_p1, "PRICING FRAMEWORK · PRINCIPLES", 0.2)
    _add_title(slide_p1, "Pricing core principles", 0.38, 0.42, 30)
    sub_p1 = slide_p1.shapes.add_textbox(Inches(0.55), Inches(0.8), Inches(12.15), Inches(0.95))
    sp1 = sub_p1.text_frame.paragraphs[0]
    sp1.text = (
        "North star: physical GMV and PDV revenue — keep the bill legible, list anchors honest to the market, "
        "and settlement + bundle economics visible where deals get made."
    )
    sp1.font.size = Pt(PT_PRINCIPLES_SUB)
    sp1.font.color.rgb = MUTED
    _font_body(sp1)
    sp1.line_spacing = 1.22

    gap_x = 0.1
    inner_w = 12.333  # 13.333 - 0.5 - 0.5
    left_m = 0.5
    w4 = (inner_w - 3 * gap_x) / 4
    w3 = (inner_w - 2 * gap_x) / 3
    # Taller cards + lower start so 12/11pt body fills the slide (less empty band below)
    y1 = 1.78
    card_h = 2.12
    row_gap = 0.1
    y2 = y1 + card_h + row_gap

    for i in range(4):
        num, icon, title, why = principles[i]
        x = left_m + i * (w4 + gap_x)
        _principle_card(
            slide_p1,
            x,
            y1,
            w4,
            card_h,
            num,
            icon,
            title,
            why,
            title_pt=PT_PRIN_COMPACT_HEAD,
            body_pt=PT_PRIN_COMPACT_BODY,
            line_before_body=7,
        )
    for i in range(3):
        num, icon, title, why = principles[4 + i]
        x = left_m + i * (w3 + gap_x)
        _principle_card(
            slide_p1,
            x,
            y2,
            w3,
            card_h,
            num,
            icon,
            title,
            why,
            title_pt=PT_PRIN_COMPACT_HEAD,
            body_pt=PT_PRIN_COMPACT_BODY,
            line_before_body=7,
        )

    src = slide_p1.shapes.add_textbox(Inches(0.55), Inches(6.88), Inches(12.0), Inches(0.42))
    sp = src.text_frame.paragraphs[0]
    sp.text = "Source: PRICING-PRINCIPLES-AND-SCENARIOS (working draft)."
    sp.font.size = Pt(11)
    sp.font.color.rgb = MUTED
    _font_body(sp)

    # --- Discussion ---
    slide_d = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide_d, BG)
    _add_kicker(slide_d, "PRICING FRAMEWORK · STEERING DISCUSSION", 0.26)
    _add_title(slide_d, "Where should we put the balance?", 0.46, 0.5, 26)
    qsub = slide_d.shapes.add_textbox(Inches(0.52), Inches(0.98), Inches(12.3), Inches(0.88))
    qp = qsub.text_frame.paragraphs[0]
    qp.text = (
        "Simplicity vs economic precision: hitting NOTR/segment targets perfectly often means more plans and exceptions. "
        "Examples: CPT waiver; invoicing inside CPT; how many offline/PdV plan variants we ship."
    )
    qp.font.size = Pt(PT_SUBTITLE)
    qp.font.color.rgb = WHITE
    _font_body(qp)
    qp.line_spacing = 1.2

    col_w = 5.95
    gap = 0.48
    lx, rx = 0.52, 0.52 + col_w + gap
    y0 = 1.88
    box_h = 1.52
    pad = 0.08

    def _boxed_text(left: float, top: float, w: float, h: float, header: str, body: str) -> None:
        shp = slide_d.shapes.add_shape(1, Inches(left), Inches(top), Inches(w), Inches(h))
        shp.fill.solid()
        shp.fill.fore_color.rgb = TBL_ROW_A
        shp.line.color.rgb = ACCENT
        shp.line.width = Pt(1.25)
        tb = slide_d.shapes.add_textbox(Inches(left + 0.12), Inches(top + 0.1), Inches(w - 0.24), Inches(h - 0.16))
        tf = tb.text_frame
        tf.word_wrap = True
        p0 = tf.paragraphs[0]
        p0.text = header
        p0.font.size = Pt(PT_BOX_HEAD)
        p0.font.bold = True
        p0.font.color.rgb = ACCENT
        _font_title(p0)
        p1 = tf.add_paragraph()
        p1.text = body
        p1.font.size = Pt(PT_BOX_BODY)
        p1.font.color.rgb = WHITE
        _font_body(p1)
        p1.space_before = Pt(6)
        p1.line_spacing = 1.22

    sim_desc = (
        "CPT waiver when merchants adopt Nuvem Pago; fold invoicing into CPT instead of a separate invoice line. "
        "Also: number of distinct offline/PdV commercial plans — each plan is another story for sales, support, and the statement."
    )
    sim_pros = "Faster merchant comprehension; easier enablement; lower ops load on edge cases."
    sim_cons = "May blur who subsidizes whom; harder to micro-tune NOTR per segment; risk of leaving revenue on the table."
    eco_desc = (
        "North star: offline GMV & revenue. We still want NOTR guardrails, non-promo list anchors, settlement visible in quoting, "
        "and a repricing loop — but perfect economics on every segment often requires more granularity than a single simple facade."
    )
    eco_pros = "Tighter BC & acquirer alignment; room for Upmarket vs SMB; clearer internal P&L per lever."
    eco_cons = "Higher cognitive load; support/sales friction; quote-cycle length; complexity debt on every new lever."

    _boxed_text(lx, y0, col_w, box_h, "SIMPLICITY — DESCRIPTION", sim_desc)
    _boxed_text(lx, y0 + box_h + pad, col_w, box_h, "SIMPLICITY — PROS", sim_pros)
    _boxed_text(lx, y0 + 2 * (box_h + pad), col_w, box_h, "SIMPLICITY — CONS", sim_cons)
    _boxed_text(rx, y0, col_w, box_h, "ECONOMIC OUTPUT — DESCRIPTION", eco_desc)
    _boxed_text(rx, y0 + box_h + pad, col_w, box_h, "ECONOMIC OUTPUT — PROS", eco_pros)
    _boxed_text(rx, y0 + 2 * (box_h + pad), col_w, box_h, "ECONOMIC OUTPUT — CONS", eco_cons)

    qfoot = slide_d.shapes.add_textbox(Inches(0.52), Inches(6.9), Inches(12.3), Inches(0.5))
    qfp = qfoot.text_frame.paragraphs[0]
    qfp.text = (
        "Committee question: maximum simplicity vs maximum precision vs hybrid "
        "(simple facade + controlled back-end segmentation) for 2026?"
    )
    qfp.font.size = Pt(PT_FOOT)
    qfp.font.color.rgb = ACCENT
    qfp.font.bold = True
    _font_body(qfp)

    prs.save(OUT)
    print(f"Wrote {OUT}")


if __name__ == "__main__":
    build()
